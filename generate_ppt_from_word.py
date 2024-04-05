import tkinter as tk
from tkinter import filedialog
from docx import Document
from pptx import Presentation
from pptx.util import Pt

def prompt_user_for_word_document():
    root = tk.Tk()
    root.withdraw()
    file_path = filedialog.askopenfilename(filetypes=[("Word documents", "*.docx")])
    return file_path

def prompt_user_for_title_and_subtitle():
    title = input("Enter the title of the presentation: ")
    subtitle = input("Enter the subtitle of the presentation: ")
    return title, subtitle

def process_word_document(file_path):
    doc = Document(file_path)
    slides_data = []
    title = None
    content = ""
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():  # If paragraph is not empty
            if title is None:
                title = paragraph.text.strip()
            else:
                content += paragraph.text.strip() + "\n\n"
        else:  # Empty paragraph indicates a new slide
            slides_data.append((title, content))
            title = None
            content = ""
    # Append the last slide
    if title is not None:
        slides_data.append((title, content))
    return slides_data

def set_font_size(shape, font_size):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

def create_presentation(title, subtitle, slides_data):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    for slide_title, slide_content in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        title_shape.text = slide_title
        set_font_size(title_shape, 40)
        content_shape.text = slide_content
        set_font_size(content_shape, 28)

    prs.save(f"{title}.pptx")

def main():
    word_doc_path = prompt_user_for_word_document()
    if not word_doc_path:
        print("No file selected. Exiting.")
        return

    title, subtitle = prompt_user_for_title_and_subtitle()
    slides_data = process_word_document(word_doc_path)
    create_presentation(title, subtitle, slides_data)
    print("Presentation created successfully.")

if __name__ == "__main__":
    main()
